"""
Générateur de présentation PowerPoint pour les résultats KNN sur le dataset CSTH.

Ce script crée une présentation professionnelle en français détaillant l'analyse
de détection de défauts sur le système Continuous Stirred Tank Heater.

Usage:
    python generate_presentation.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import pandas as pd
import json


# ============================================================================
# Configuration des couleurs et styles
# ============================================================================

# Palette de couleurs professionnelle (bleu/gris)
COLOR_PRIMARY = RGBColor(31, 78, 121)      # Bleu foncé
COLOR_SECONDARY = RGBColor(68, 114, 196)   # Bleu moyen
COLOR_ACCENT = RGBColor(237, 125, 49)      # Orange
COLOR_SUCCESS = RGBColor(112, 173, 71)     # Vert
COLOR_WARNING = RGBColor(255, 192, 0)      # Jaune
COLOR_DANGER = RGBColor(192, 0, 0)         # Rouge
COLOR_TEXT = RGBColor(64, 64, 64)          # Gris foncé
COLOR_LIGHT_GRAY = RGBColor(217, 217, 217) # Gris clair
COLOR_BG = RGBColor(242, 242, 242)         # Fond gris très clair


# ============================================================================
# Fonctions utilitaires pour le formatage
# ============================================================================

def add_title_slide(prs, title, subtitle=""):
    """Ajoute une diapositive de titre."""
    slide_layout = prs.slide_layouts[0]  # Layout titre
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    if subtitle:
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT
    
    return slide


def add_section_slide(prs, section_title):
    """Ajoute une diapositive de section."""
    slide_layout = prs.slide_layouts[6]  # Layout vide
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond coloré
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Titre centré
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = section_title
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(54)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    
    return slide


def add_content_slide(prs, title):
    """Ajoute une diapositive de contenu avec titre."""
    slide_layout = prs.slide_layouts[5]  # Layout vide
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    paragraph = title_frame.paragraphs[0]
    paragraph.font.size = Pt(32)
    paragraph.font.bold = True
    paragraph.font.color.rgb = COLOR_PRIMARY
    
    # Ligne de séparation
    left = Inches(0.5)
    top = Inches(1.1)
    width = Inches(9)
    height = Inches(0)
    
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = COLOR_SECONDARY
    line.line.width = Pt(2)
    
    return slide


def add_text_box(slide, left, top, width, height, text, font_size=14, 
                 bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT):
    """Ajoute une zone de texte formatée."""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    
    return text_box


def add_bullet_points(slide, left, top, width, height, points, font_size=16):
    """Ajoute une liste à puces."""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6)
    
    return text_box


def add_table(slide, left, top, data, col_widths=None, header_color=COLOR_PRIMARY):
    """Ajoute un tableau formaté."""
    rows = len(data)
    cols = len(data[0])
    
    # Largeurs par défaut
    if col_widths is None:
        col_widths = [Inches(2)] * cols
    
    # Créer le tableau
    table = slide.shapes.add_table(rows, cols, left, top, 
                                    sum(col_widths), Inches(0.4) * rows).table
    
    # Définir les largeurs de colonnes
    for i, width in enumerate(col_widths):
        table.columns[i].width = width
    
    # Remplir le tableau
    for i, row_data in enumerate(data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_data)
            
            # Style de l'en-tête
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(14)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Style des cellules
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
                
                # Alternance de couleurs
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
    
    return table


def add_metric_box(slide, left, top, width, height, label, value, color=COLOR_SECONDARY):
    """Ajoute une boîte de métrique stylisée."""
    # Fond
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    
    # Valeur
    value_box = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(0.6))
    value_frame = value_box.text_frame
    value_frame.text = str(value)
    p = value_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Label
    label_box = slide.shapes.add_textbox(left, top + Inches(0.85), width, Inches(0.3))
    label_frame = label_box.text_frame
    label_frame.text = label
    p = label_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)


# ============================================================================
# Fonctions de création des diapositives
# ============================================================================

def create_title_slide(prs):
    """Diapositive de titre."""
    add_title_slide(
        prs,
        "Classification KNN pour la Détection de Défauts",
        "Système de Chauffage à Réservoir Agité (CSTH)"
    )


def create_agenda_slide(prs):
    """Diapositive agenda."""
    slide = add_content_slide(prs, "Plan de Présentation")
    
    points = [
        "1. Introduction et Contexte",
        "2. Dataset CSTH",
        "3. Méthodologie KNN",
        "4. Recherche d'Hyperparamètres",
        "5. Résultats Finaux",
        "6. Analyse et Conclusions"
    ]
    
    add_bullet_points(slide, Inches(1.5), Inches(2), Inches(7), Inches(4), points, 22)


def create_context_slide(prs):
    """Diapositive de contexte."""
    slide = add_content_slide(prs, "Contexte : Détection de Défauts Industriels")
    
    points = [
        "Objectif : Détecter automatiquement les défauts instrumentaux dans un système de chauffage",
        "Application : Maintenance prédictive et surveillance en temps réel",
        "Approche : Classification supervisée par K-plus proches voisins (KNN)",
        "Données : Séries temporelles simulées du système CSTH"
    ]
    
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(3.5), points, 18)


def create_dataset_slide(prs):
    """Diapositive description du dataset."""
    slide = add_content_slide(prs, "Dataset CSTH - Caractéristiques")
    
    # Description
    desc_text = (
        "Continuous Stirred Tank Heater : système de chauffage avec mélange "
        "d'eau chaude/froide, chauffage vapeur et contrôle en boucle fermée"
    )
    add_text_box(slide, Inches(1), Inches(1.8), Inches(8), Inches(0.8), 
                 desc_text, 16, color=COLOR_TEXT)
    
    # Tableau des caractéristiques
    data = [
        ["Caractéristique", "Valeur"],
        ["Échantillons d'entraînement", "6 300"],
        ["Échantillons de validation", "900"],
        ["Échantillons de test", "1 800"],
        ["Dimension temporelle", "3 pas de temps"],
        ["Nombre de features", "200"],
        ["Classes", "2 (normal / défaut)"],
        ["Distribution", "~50% / ~50%"]
    ]
    
    add_table(slide, Inches(2), Inches(3), data, 
              [Inches(3), Inches(2)])


def create_methodology_slide(prs):
    """Diapositive méthodologie."""
    slide = add_content_slide(prs, "Méthodologie : Pipeline KNN")
    
    # Étapes
    steps = [
        "1. Prétraitement",
        "   • Standardisation des données (moyenne=0, écart-type=1)",
        "   • Réduction de dimensionnalité par ACP (95% variance)",
        "",
        "2. Calcul des distances",
        "   • Distance euclidienne sur séries temporelles aplaties",
        "",
        "3. Classification KNN",
        "   • Vote pondéré par distance inverse",
        "   • Validation croisée par groupes (5 folds)"
    ]
    
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4), steps, 16)


def create_hyperparam_slide(prs):
    """Diapositive recherche d'hyperparamètres."""
    slide = add_content_slide(prs, "Recherche d'Hyperparamètres : Valeur de k")
    
    # Résultats
    data = [
        ["k", "Accuracy", "F1-Score", "Temps (s)"],
        ["1", "0.729", "0.728", "-"],
        ["3", "0.743", "0.743", "-"],
        ["5", "0.743", "0.743", "-"],
        ["7", "0.762", "0.762", "-"],
        ["10", "0.786", "0.784", "-"],
        ["15", "0.790", "0.787", "-"],
        ["20", "0.787", "0.783", "-"],
        ["25", "0.796", "0.792", "-"],
        ["30", "0.793", "0.788", "-"]
    ]
    
    add_table(slide, Inches(2.5), Inches(2.2), data,
              [Inches(0.8), Inches(1.2), Inches(1.2), Inches(1.2)])
    
    # Conclusion
    add_text_box(slide, Inches(1), Inches(5.8), Inches(8), Inches(0.5),
                 "→ Meilleure performance : k = 25 (F1 = 0.792)",
                 18, bold=True, color=COLOR_ACCENT)


def create_results_slide(prs):
    """Diapositive résultats finaux."""
    slide = add_content_slide(prs, "Résultats Finaux sur Ensemble de Test")
    
    # Métriques principales
    add_metric_box(slide, Inches(0.8), Inches(2), Inches(2.2), Inches(1.3),
                   "Accuracy", "83.61%", COLOR_SUCCESS)
    add_metric_box(slide, Inches(3.4), Inches(2), Inches(2.2), Inches(1.3),
                   "F1-Score", "83.45%", COLOR_SECONDARY)
    add_metric_box(slide, Inches(6), Inches(2), Inches(2.2), Inches(1.3),
                   "Temps", "0.70s", COLOR_PRIMARY)
    
    # Détails par classe
    add_text_box(slide, Inches(1), Inches(3.8), Inches(8), Inches(0.4),
                 "Performance par Classe", 20, bold=True, color=COLOR_PRIMARY)
    
    data = [
        ["Classe", "Précision", "Rappel", "F1-Score", "Support"],
        ["Normal", "0.917", "0.738", "0.818", "897"],
        ["Défaut", "0.782", "0.934", "0.851", "903"]
    ]
    
    add_table(slide, Inches(1.5), Inches(4.3), data,
              [Inches(1.5), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2)])


def create_confusion_matrix_slide(prs):
    """Diapositive matrice de confusion."""
    slide = add_content_slide(prs, "Matrice de Confusion")
    
    # Titre de la matrice
    add_text_box(slide, Inches(2), Inches(2), Inches(6), Inches(0.4),
                 "Prédictions vs Réalité", 18, bold=True, 
                 color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
    
    # Matrice
    data = [
        ["", "Prédit: Normal", "Prédit: Défaut"],
        ["Réel: Normal", "662 (VN)", "235 (FP)"],
        ["Réel: Défaut", "60 (FN)", "843 (VP)"]
    ]
    
    table = add_table(slide, Inches(2.5), Inches(2.8), data,
                      [Inches(1.5), Inches(1.5), Inches(1.5)],
                      header_color=COLOR_PRIMARY)
    
    # Colorer les cellules
    table.cell(1, 1).fill.solid()
    table.cell(1, 1).fill.fore_color.rgb = RGBColor(198, 224, 180)
    
    table.cell(2, 2).fill.solid()
    table.cell(2, 2).fill.fore_color.rgb = RGBColor(198, 224, 180)
    
    table.cell(1, 2).fill.solid()
    table.cell(1, 2).fill.fore_color.rgb = RGBColor(255, 230, 153)
    
    table.cell(2, 1).fill.solid()
    table.cell(2, 1).fill.fore_color.rgb = RGBColor(255, 230, 153)
    
    # Légende
    legend = [
        "VN (Vrai Négatif) : Normal correctement identifié",
        "VP (Vrai Positif) : Défaut correctement détecté",
        "FP (Faux Positif) : Fausse alarme (26.2%)",
        "FN (Faux Négatif) : Défaut manqué (6.6%)"
    ]
    
    add_bullet_points(slide, Inches(1), Inches(5.2), Inches(8), Inches(2), legend, 14)


def create_fault_metrics_slide(prs):
    """Diapositive métriques de détection de défauts."""
    slide = add_content_slide(prs, "Métriques Spécifiques à la Détection de Défauts")
    
    # Métriques
    add_metric_box(slide, Inches(1.5), Inches(2.2), Inches(3), Inches(1.5),
                   "Taux de Détection\n(Rappel)", "93.36%", COLOR_SUCCESS)
    add_metric_box(slide, Inches(5.5), Inches(2.2), Inches(3), Inches(1.5),
                   "Taux de\nFausses Alarmes", "26.20%", COLOR_WARNING)
    
    # Analyse
    analysis = [
        "Points forts :",
        "  • Excellente détection des défauts (93.4% de rappel)",
        "  • Peu de défauts manqués (60/903 = 6.6%)",
        "",
        "Points d'amélioration :",
        "  • Taux de fausses alarmes élevé (26.2%)",
        "  • 235 échantillons normaux classés comme défauts",
        "",
        "Compromis performance-coût :",
        "  • Configuration adaptée aux cas où manquer un défaut coûte cher",
        "  • Acceptable si vérification humaine post-alarme possible"
    ]
    
    add_bullet_points(slide, Inches(1), Inches(4.2), Inches(8), Inches(2.5), analysis, 15)


def create_preprocessing_slide(prs):
    """Diapositive impact du prétraitement."""
    slide = add_content_slide(prs, "Impact du Prétraitement (Étude d'Ablation)")
    
    add_text_box(slide, Inches(1), Inches(1.8), Inches(8), Inches(0.5),
                 "Comparaison des différentes configurations de pipeline",
                 16, color=COLOR_TEXT)
    
    # Résultats
    data = [
        ["Configuration", "Accuracy", "F1-weighted"],
        ["Sans prétraitement", "~0.73", "~0.73"],
        ["Standardisation seule", "~0.76", "~0.76"],
        ["ACP seule (95% var)", "~0.75", "~0.75"],
        ["Pipeline complet (std+ACP)", "0.796", "0.791"]
    ]
    
    add_table(slide, Inches(1.5), Inches(2.8), data,
              [Inches(3.5), Inches(1.5), Inches(1.5)])
    
    # Conclusion
    conclusion = [
        "→ Le pipeline complet (standardisation + ACP) offre les meilleures performances",
        "→ Gain de ~6 points de pourcentage par rapport à la baseline",
        "→ Réduction de dimensionnalité (600 → 6 features) tout en préservant 95% de variance"
    ]
    
    add_bullet_points(slide, Inches(1), Inches(5.2), Inches(8), Inches(2), conclusion, 15)


def create_conclusions_slide(prs):
    """Diapositive conclusions."""
    slide = add_content_slide(prs, "Conclusions et Perspectives")
    
    # Conclusions
    conclusions = [
        "Résultats obtenus :",
        "  • 83.6% de précision globale sur données de test",
        "  • Excellent taux de détection des défauts (93.4%)",
        "  • Pipeline efficace et rapide (<1s pour 1800 échantillons)",
        "",
        "Forces de l'approche KNN :",
        "  • Simplicité et interprétabilité",
        "  • Pas de phase d'entraînement coûteuse",
        "  • Performance robuste avec prétraitement adéquat",
        "",
        "Pistes d'amélioration :",
        "  • Réduire le taux de fausses alarmes (ajustement du seuil)",
        "  • Tester DTW pour mieux capturer la dynamique temporelle",
        "  • Explorer des ensembles de modèles (KNN + autres)"
    ]
    
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), conclusions, 16)


def create_end_slide(prs):
    """Diapositive de fin."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Texte
    add_text_box(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                 "Merci de votre attention", 44, bold=True,
                 color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(2), Inches(4.2), Inches(6), Inches(0.5),
                 "Questions ?", 28,
                 color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)


# ============================================================================
# Fonction principale
# ============================================================================

def generate_presentation(output_path="presentation_knn_csth.pptx"):
    """Génère la présentation complète."""
    
    print("Génération de la présentation PowerPoint...")
    
    # Créer la présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Diapositives
    print("  [1/12] Diapositive de titre...")
    create_title_slide(prs)
    
    print("  [2/12] Agenda...")
    create_agenda_slide(prs)
    
    print("  [3/12] Section Introduction...")
    add_section_slide(prs, "Introduction")
    
    print("  [4/12] Contexte...")
    create_context_slide(prs)
    
    print("  [5/12] Dataset...")
    create_dataset_slide(prs)
    
    print("  [6/12] Section Méthodologie...")
    add_section_slide(prs, "Méthodologie")
    
    print("  [7/12] Pipeline KNN...")
    create_methodology_slide(prs)
    
    print("  [8/12] Section Résultats...")
    add_section_slide(prs, "Résultats")
    
    print("  [9/12] Hyperparamètres...")
    create_hyperparam_slide(prs)
    
    print("  [10/12] Résultats finaux...")
    create_results_slide(prs)
    
    print("  [11/12] Matrice de confusion...")
    create_confusion_matrix_slide(prs)
    
    print("  [12/12] Métriques de détection...")
    create_fault_metrics_slide(prs)
    
    print("  [13/13] Prétraitement...")
    create_preprocessing_slide(prs)
    
    print("  [14/14] Conclusions...")
    create_conclusions_slide(prs)
    
    print("  [15/15] Diapositive de fin...")
    create_end_slide(prs)
    
    # Sauvegarder
    prs.save(output_path)
    print(f"\n✓ Présentation générée avec succès : {output_path}")
    print(f"  Nombre de diapositives : {len(prs.slides)}")


if __name__ == "__main__":
    generate_presentation()
